import os
import io
from flask import Flask, request, jsonify, Response, stream_with_context, send_from_directory
from flask_cors import CORS
import anthropic
import docx as docx_lib
import pdfplumber
from pptx import Presentation
import openpyxl

app = Flask(__name__)
_origins = [o.strip() for o in os.environ.get("ALLOWED_ORIGINS", "http://localhost:5557,http://localhost:3000").split(",")]
CORS(app, origins=_origins)
app.config["MAX_CONTENT_LENGTH"] = 20 * 1024 * 1024  # 20 MB

_key = os.environ.get("ANTHROPIC_API_KEY", "")
if not _key:
    raise RuntimeError("ANTHROPIC_API_KEY environment variable is not set")
client = anthropic.Anthropic(api_key=_key)

MAX_CHARS = 120_000
ALLOWED_EXT = {".docx", ".pdf", ".pptx", ".xlsx"}

# ══════════════════════════════════════════════════════════════════════════════
#  SYSTEM PROMPTS
# ══════════════════════════════════════════════════════════════════════════════

SUPERVISOR_PROMPT = """You are a senior academic supervisor reviewing a student's work. You are invested in their growth and development as a thinker and writer. You are strict about academic standards and intellectual rigour — you won't accept vague claims, unsupported assertions, or lazy structure — but your tone is engaged and mentoring, like sitting across a desk working through the paper together. Every comment you make teaches something; you never just flag a problem without explaining why it matters.

Good work deserves recognition. You actively look for and call out what the student has done well — a well-constructed argument, effective use of evidence, clear structure, precise language. Positive feedback is not filler; it tells the student what to keep and build on.

You review work across seven dimensions:
- strength: something genuinely done well — a strong argument, effective structure, good evidence use, clear writing
- grammar: spelling, punctuation, syntax
- content: argument quality, logical coherence, depth of analysis
- structure: organisation, paragraph construction, flow between ideas, introduction/conclusion quality
- clarity: ambiguous phrasing, wordiness, passive voice overuse, readability
- evidence: use of data, sourcing, statistical claims, quality of supporting material
- citations: APA 7 compliance — in-text format (Author, Year), reference list ordering, DOIs, author name formatting

CRITICAL RULE ON SUGGESTIONS: Your suggestion field must tell the student HOW to think about improving this — the approach, the structure, what to consider. You must NEVER write replacement prose. You are guiding their thinking, not writing for them. For strength comments, the suggestion can note how to apply this strength elsewhere in the paper, or leave it null.

Output format: NDJSON — one JSON object per line, nothing else. No preamble, no markdown, no code fences.

Each comment:
{"paragraph_index": 2, "type": "content", "quote": "verbatim phrase from text ≤15 words", "comment": "what is weak or strong and why it matters", "suggestion": "how the student should approach fixing or extending this — null for praise-only strengths", "severity": "low|medium|high"}

After all comments, one final line:
{"type": "summary", "overall_grade": null, "summary_text": "Conversational 2-3 sentence summary: what is genuinely working well, the single most important thing to fix before the next draft, and one concrete action to take."}

Rules:
- paragraph_index must be an integer matching a [N] index present in the document
- quote must be verbatim from that paragraph, ≤15 words
- type is exactly one of: strength, grammar, content, structure, clarity, evidence, citations
- strength comments must be genuine — do not manufacture praise; only flag what is actually good
- include at least 3–5 strength comments across the document, spread naturally throughout
- severity for strength comments is always "low" (it is not a problem scale, just required for schema consistency)
- severity for issue comments: low = minor polish, medium = noticeably weakens the work, high = fundamental problem
- emit 15–50 comments depending on document length; aim for roughly 25–35% strength, 65–75% issues
- CRITICAL: distribute comments across the ENTIRE document — your paragraph_index values must span from the opening paragraphs to the final paragraphs; do not cluster in the first half
- overall_grade is always null in supervisor mode
- never write a generic comment — name the specific thing and explain why it matters"""

GRADING_PROMPT = """You are a university professor formally evaluating a student's submitted work. You are an assessor, not a mentor. Your comments read like written marginal notes on a returned paper. You are thorough, you do not skip minor issues, and your APA 7 enforcement is strict.

A fair assessment acknowledges what the student has done well, not just where they have failed. You note genuine strengths — a well-structured argument, effective use of evidence, clear writing — because accurate grading requires recognising both achievement and shortcoming.

You evaluate across seven dimensions:
- strength: something genuinely done well — a strong argument, clear structure, effective evidence, precise language
- grammar: spelling, punctuation, syntax
- content: argument quality, logical coherence, analytical depth, factual accuracy
- structure: organisation, paragraph construction, flow, introduction/conclusion quality
- clarity: ambiguous phrasing, wordiness, passive voice, readability
- evidence: use of data, sourcing, statistical claims, missing or weak support
- citations: APA 7 — in-text format (Author, Year), reference list alphabetical order, DOIs required where available, author name formatting (Last, F. M.), hanging indents, et al. usage

If a rubric is provided, evaluate the work against it explicitly. Reference rubric criteria in your comments where relevant.

CRITICAL RULE ON SUGGESTIONS: Your suggestion must tell the student what they need to do conceptually — what kind of evidence to find, what structural change to make, what the argument is missing. You must NEVER write replacement prose or dictate exact wording. For strength comments, suggestion can note how to carry this quality through the rest of the work, or be null.

Output format: NDJSON — one JSON object per line, nothing else. No preamble, no markdown, no code fences.

Each comment:
{"paragraph_index": 2, "type": "citations", "quote": "verbatim phrase from text ≤15 words", "comment": "precise identification of the strength or problem and its academic impact", "suggestion": "what the student needs to do — null for praise-only strengths", "severity": "low|medium|high"}

After all comments, one final line:
{"type": "summary", "overall_grade": "B-", "summary_text": "2-4 sentence formal assessment: the submission's main strength, its most critical academic failure, and whether it meets the standard for this level of study."}

Rules:
- paragraph_index must be an integer matching a [N] index present in the document
- quote must be verbatim from that paragraph, ≤15 words
- type is exactly one of: strength, grammar, content, structure, clarity, evidence, citations
- strength comments must be genuine — do not manufacture praise
- include at least 2–4 strength comments spread across the document
- severity for strength comments is always "low"
- severity for issue comments: low = minor, medium = weakens the grade, high = submission-level failure
- emit 15–55 comments; do not skip minor issues in grading mode
- CRITICAL: distribute comments across the ENTIRE document — your paragraph_index values must span from the opening paragraphs to the final paragraphs; do not cluster in the first half
- overall_grade uses A+, A, A-, B+, B, B-, C+, C, C-, D, F
- never write a generic comment"""

CHAT_PROMPT = """You are a senior academic supervisor having a conversation with a student about their paper. You have already read the document (it is in this conversation). You are engaged, direct, and genuinely interested in helping them improve their thinking and writing. You explain things clearly, you ask good questions back when it helps them think, and you never just hand them the answer — you guide them toward it.

You do NOT write prose for the student. If they ask you to rewrite a sentence or paragraph, explain what is wrong with it and how they should think about fixing it, but do not produce a corrected version. Your role is to develop their ability, not substitute for it.

Respond in plain conversational prose — no bullet lists unless the student asks for one, no NDJSON, no markdown headers. Keep responses focused and useful, not lengthy."""

# ══════════════════════════════════════════════════════════════════════════════
#  FILE PARSING
# ══════════════════════════════════════════════════════════════════════════════

def _merge_pdf_lines(lines):
    """Rejoin PDF lines that were split mid-sentence."""
    merged, buf = [], ""
    for line in lines:
        if not line.strip():
            if buf:
                merged.append(buf.strip())
                buf = ""
            continue
        if buf and not buf[-1] in ".?!:":
            buf += " " + line.strip()
        else:
            if buf:
                merged.append(buf.strip())
            buf = line.strip()
    if buf:
        merged.append(buf.strip())
    return [m for m in merged if m]


def parse_docx(file_bytes):
    doc = docx_lib.Document(io.BytesIO(file_bytes))
    paragraphs = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            paragraphs.append(t)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))
    return paragraphs


def parse_pdf(file_bytes):
    lines = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            text = page.extract_text(x_tolerance=2, y_tolerance=3)
            if text:
                for line in text.split("\n"):
                    lines.append(line.strip())
    return _merge_pdf_lines(lines)


def parse_pptx(file_bytes):
    prs = Presentation(io.BytesIO(file_bytes))
    paragraphs = []
    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        if texts:
            paragraphs.append(f"[Slide {slide_num}] " + " | ".join(texts))
    return paragraphs


def parse_xlsx(file_bytes):
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    try:
        paragraphs = []
        for sheet in wb.worksheets:
            paragraphs.append(f"[Sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None and str(c).strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))
        return paragraphs
    finally:
        wb.close()


def extract_text(filename, file_bytes):
    ext = os.path.splitext(filename)[1].lower()
    if ext == ".docx":  return parse_docx(file_bytes)
    if ext == ".pdf":   return parse_pdf(file_bytes)
    if ext == ".pptx":  return parse_pptx(file_bytes)
    if ext == ".xlsx":  return parse_xlsx(file_bytes)
    raise ValueError(f"Unsupported file type: {ext}")


def chunk_paragraphs(paragraphs, max_chars=MAX_CHARS):
    result, total = [], 0
    for p in paragraphs:
        if total + len(p) > max_chars:
            break
        result.append(p)
        total += len(p) + 2
    return result, len(result) < len(paragraphs)


def numbered_doc(paragraphs):
    return "\n\n".join(f"[{i}] {p}" for i, p in enumerate(paragraphs) if p.strip())

# ══════════════════════════════════════════════════════════════════════════════
#  STREAMING GENERATORS
# ══════════════════════════════════════════════════════════════════════════════

def _build_context_block(context):
    if not context:
        return ""
    lines = []
    if context.get("subject"): lines.append(f"Subject/Course: {context['subject']}")
    if context.get("type"):    lines.append(f"Submission type: {context['type']}")
    if context.get("topic"):   lines.append(f"Topic/Title: {context['topic']}")
    if context.get("notes"):   lines.append(f"Additional context: {context['notes']}")
    if not lines:
        return ""
    return "\n\n[SUBMISSION CONTEXT]\n" + "\n".join(lines) + "\n[/SUBMISSION CONTEXT]"


def _review_stream(paragraphs, mode, rubric="", context=None):
    prompt = GRADING_PROMPT if mode == "grading" else SUPERVISOR_PROMPT
    doc = numbered_doc(paragraphs)
    ctx_block    = _build_context_block(context)
    rubric_block = f"\n\n[RUBRIC]\n{rubric.strip()}\n[/RUBRIC]" if rubric and rubric.strip() else ""
    user_content = f"Please review the following document.{ctx_block}{rubric_block}\n\n{doc}"
    with client.messages.stream(
        model="claude-sonnet-4-6",
        max_tokens=8192,
        system=[{"type": "text", "text": prompt, "cache_control": {"type": "ephemeral"}}],
        messages=[{"role": "user", "content": user_content}],
    ) as stream:
        for text in stream.text_stream:
            yield text


def _chat_stream(messages):
    with client.messages.stream(
        model="claude-sonnet-4-6",
        max_tokens=2048,
        system=[{"type": "text", "text": CHAT_PROMPT, "cache_control": {"type": "ephemeral"}}],
        messages=messages,
    ) as stream:
        for text in stream.text_stream:
            yield text

# ══════════════════════════════════════════════════════════════════════════════
#  ROUTES
# ══════════════════════════════════════════════════════════════════════════════

@app.route("/")
def index():
    return send_from_directory(".", "index.html")


@app.route("/api/upload", methods=["POST"])
def upload():
    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400

    f = request.files["file"]
    filename = f.filename or ""
    ext = os.path.splitext(filename)[1].lower()

    if ext not in ALLOWED_EXT:
        return jsonify({"error": f"Unsupported file type. Please upload .docx, .pdf, .pptx, or .xlsx"}), 400

    try:
        file_bytes = f.read()
    except Exception:
        return jsonify({"error": "Could not read file"}), 422

    try:
        paragraphs = extract_text(filename, file_bytes)
    except Exception as e:
        return jsonify({"error": f"Could not parse file: {str(e)}"}), 422

    if not paragraphs:
        return jsonify({"error": "No readable text found in this file."}), 422

    trimmed, truncated = chunk_paragraphs(paragraphs)
    return jsonify({
        "paragraphs": trimmed,
        "truncated": truncated,
        "char_count": sum(len(p) for p in trimmed)
    })


@app.route("/api/review", methods=["POST"])
def review():
    data = request.get_json(silent=True) or {}
    paragraphs = data.get("paragraphs", [])
    mode    = data.get("mode", "supervisor")
    rubric  = data.get("rubric", "")
    context = data.get("context") or {}

    if not paragraphs:
        return jsonify({"error": "No paragraphs provided"}), 400
    if mode not in ("supervisor", "grading"):
        mode = "supervisor"

    try:
        return Response(
            stream_with_context(_review_stream(paragraphs, mode, rubric, context)),
            content_type="text/plain;charset=utf-8"
        )
    except Exception as e:
        return jsonify({"error": f"AI service error: {str(e)}"}), 502


@app.route("/api/chat", methods=["POST"])
def chat():
    data = request.get_json(silent=True) or {}
    paragraphs = data.get("paragraphs", [])
    history    = data.get("history", [])
    user_msg   = (data.get("message") or "").strip()
    context    = data.get("context") or {}

    if not user_msg:
        return jsonify({"error": "No message provided"}), 400

    if history:
        messages = history + [{"role": "user", "content": user_msg}]
    else:
        if not paragraphs:
            return jsonify({"error": "No document loaded"}), 400
        doc = numbered_doc(paragraphs)
        ctx_block = _build_context_block(context)
        intro = f"Here is the document we'll be discussing.{ctx_block}\n\n{doc}"
        messages = [{
            "role": "user",
            "content": [
                {"type": "text", "text": intro, "cache_control": {"type": "ephemeral"}},
                {"type": "text", "text": user_msg}
            ]
        }]

    try:
        return Response(
            stream_with_context(_chat_stream(messages)),
            content_type="text/plain;charset=utf-8"
        )
    except Exception as e:
        return jsonify({"error": f"AI service error: {str(e)}"}), 502


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5557))
    print(f"Study Buddy running → http://localhost:{port}")
    app.run(host="0.0.0.0", port=port, debug=False)
